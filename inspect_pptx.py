from pptx import Presentation
import sys
import os

def inspect_presentation(file_path):
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return

    print(f"Inspecting: {file_path}")
    prs = Presentation(file_path)
    
    print(f"Slide Width: {prs.slide_width.inches} inches")
    print(f"Slide Height: {prs.slide_height.inches} inches\n")

    print(f"Number of Slide Layouts: {len(prs.slide_layouts)}")
    print("-" * 40)
    
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\nLayout {i}: {layout.name}")
        print(f"  Placeholders: {len(layout.placeholders)}")
        for shape in layout.placeholders:
            idx = shape.placeholder_format.idx
            ptype = shape.placeholder_format.type
            name = shape.name
            print(f"    - Index: {idx} | Type: {ptype} | Name: {name}")
            try:
                print(f"      Size / Pos: L={shape.left.inches:.2f}, T={shape.top.inches:.2f}, W={shape.width.inches:.2f}, H={shape.height.inches:.2f}")
            except Exception as e:
                print(f"      Size / Pos: Could not determine ({e})")

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Inspect a Slide Master for Layouts and Placeholders.")
    parser.add_argument("file", help="Path to the .pptx file to inspect")
    
    args = parser.parse_args()
    inspect_presentation(args.file)
