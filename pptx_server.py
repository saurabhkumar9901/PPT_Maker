import os
from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Initialize FastMCP Server
mcp = FastMCP("PowerPoint MCP Server")

# Premium Settings & Margins
MARGIN_LEFT = Inches(0.5)
MARGIN_TOP = Inches(1.8)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
CONTENT_WIDTH = SLIDE_WIDTH - Inches(1.0)
CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - Inches(0.5)

# Layout Mappings
def get_theme_color(theme_str: str):
    """Maps dynamic Gemini theme strings back to native Slide Master theme enums."""
    mapping = {
        "BACKGROUND_1": MSO_THEME_COLOR.BACKGROUND_1,
        "BACKGROUND_2": MSO_THEME_COLOR.BACKGROUND_2,
        "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
        "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
        "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
        "DARK_1": MSO_THEME_COLOR.DARK_1,
        "DARK_2": MSO_THEME_COLOR.DARK_2,
        "LIGHT_1": MSO_THEME_COLOR.LIGHT_1,
    }
    return mapping.get(theme_str, MSO_THEME_COLOR.ACCENT_1)

def get_cover_layout(prs):
    """Finds the cover slide layout dynamically."""
    for layout in prs.slide_layouts:
        if "cover" in layout.name.lower(): return layout
    return prs.slide_layouts[0]

def get_title_only_layout(prs):
    """Finds the title-only layout dynamically to fix missing title bugs across templates."""
    for layout in prs.slide_layouts:
        if "title only" in layout.name.lower(): return layout
    
    # Fallback to the first layout that actually has a title placeholder
    for layout in prs.slide_layouts:
        if layout.shapes.title: return layout
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

def apply_title_styling(shape, title_text, palette):
    if not shape or not palette: return
    shape.text = title_text
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.name = "Arial"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.theme_color = get_theme_color(palette[2])
        p.alignment = PP_ALIGN.LEFT

@mcp.tool()
def create_presentation(template_path: str, output_path: str) -> str:
    """Creates a new presentation from a template file and saves it to output_path."""
    if not os.path.exists(template_path): return f"Error: Template not found at {template_path}"
    prs = Presentation(template_path)
    
    # Securely remove existing template slides while preserving the Slide Master Layouts
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for slide in slides:
        prs.part.drop_rel(slide.rId) # Safely drop relationship to avoid duplicate xml warnings
        xml_slides.remove(slide)
        
    prs.save(output_path)
    return f"Created new presentation from template at {output_path}"

@mcp.tool()
def add_title_slide(file_path: str, title: str, subtitle: str, palette: list[str]) -> str:
    """Adds a Cover/Title slide to the presentation."""
    prs = Presentation(file_path)
    layout = get_cover_layout(prs)
    slide = prs.slides.add_slide(layout)
    placeholders = list(slide.placeholders)
    if len(placeholders) >= 2:
        placeholders.sort(key=lambda x: x.top)
        t_shape = placeholders[1] # Typically the larger lower box, or depending on template
        s_shape = placeholders[0]
        
        t_shape.text = title
        if t_shape.text_frame.paragraphs:
            t_shape.text_frame.paragraphs[0].font.size = Pt(44)
            t_shape.text_frame.paragraphs[0].font.bold = True
            t_shape.text_frame.paragraphs[0].font.color.theme_color = get_theme_color(palette[2])
            
        s_shape.text = subtitle
        if s_shape.text_frame.paragraphs:
            s_shape.text_frame.paragraphs[0].font.size = Pt(24)
            s_shape.text_frame.paragraphs[0].font.color.theme_color = get_theme_color(palette[3])
    prs.save(file_path)
    return "Title slide added."

@mcp.tool()
def add_content_slide(file_path: str, title: str, paragraphs: list[str], palette: list[str]) -> str:
    """Adds a content slide using a dynamic multi-column layout instead of standard bullets."""
    prs = Presentation(file_path)
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    apply_title_styling(slide.shapes.title, title.upper(), palette)

    num_items = len(paragraphs)
    if num_items == 0:
        prs.save(file_path)
        return "Content slide added."

    # Decide layout array: Max 4 columns. If 5 items, we can do 3 top, 2 bottom.
    cols = min(num_items, 4) if num_items > 0 else 1
    rows = (num_items + cols - 1) // cols
    
    col_gap = Inches(0.4)
    row_gap = Inches(0.4)
    box_width = (CONTENT_WIDTH - (col_gap * (cols - 1))) / cols
    box_height = (CONTENT_HEIGHT - (row_gap * (rows - 1))) / rows # Fill available space aggressively

    # Add a subtle background to fill blank space
    content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), MARGIN_TOP - Inches(0.2), SLIDE_WIDTH - Inches(0.6), SLIDE_HEIGHT - MARGIN_TOP)
    content_bg.fill.solid()
    content_bg.fill.fore_color.theme_color = get_theme_color(palette[0])
    content_bg.line.fill.background() # No border

    for i, para_text in enumerate(paragraphs):
        c = i % cols
        r = i // cols
        
        # Center align the row
        items_in_row = cols
        if r == rows - 1 and num_items % cols != 0:
            items_in_row = num_items % cols
            
        row_width = (items_in_row * box_width) + ((items_in_row - 1) * col_gap)
        start_x = MARGIN_LEFT + (CONTENT_WIDTH - row_width) / 2
        
        left = start_x + (box_width + col_gap) * c
        top = MARGIN_TOP + (box_height + row_gap) * r
        
        # Add Card Background
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height)
        bg.fill.solid()
        bg.fill.fore_color.theme_color = get_theme_color(palette[1])
        bg.line.color.theme_color = get_theme_color(palette[3])
        bg.line.width = Pt(1.5)
        
        # Add Text Inside Card
        tf = bg.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Middle align to fix negative space issue
        tf.margin_top = Inches(0.3)
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_bottom = Inches(0.3)
        
        # Split into pseudo-header if it contains a colon
        parts = para_text.split(':', 1)
        
        p = tf.paragraphs[0]
        p.text = parts[0] + (':' if len(parts) > 1 else '')
        p.font.size = Pt(32) if any(char.isdigit() for char in parts[0]) else Pt(20) # Uniform sizes
        p.font.bold = True
        p.font.color.theme_color = get_theme_color(palette[3])
        
        if len(parts) > 1:
            p2 = tf.add_paragraph()
            p2.text = parts[1].strip()
            p2.font.size = Pt(14) # Uniform safely sized text
            p2.font.color.theme_color = get_theme_color(palette[2])
            p2.space_before = Pt(8)

    prs.save(file_path)
    return "Content slide added."

@mcp.tool()
def add_infographic_slide(file_path: str, title: str, steps: list[dict], palette: list[str]) -> str:
    """Adds a process flow infographic with chevrons or connected boxes."""
    prs = Presentation(file_path)
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    apply_title_styling(slide.shapes.title, title.upper(), palette)

    num_steps = len(steps)
    if num_steps == 0:
        prs.save(file_path)
        return "Infographic slide added with no steps."
        
    total_gap_width = Inches(0.6) * (num_steps - 1)
    box_width = min(Inches(3.0), (CONTENT_WIDTH - total_gap_width) / num_steps)
    box_height = Inches(2.2)
    top_offset = Inches(3.0)
    
    # Add Content background
    content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top_offset - Inches(0.6), SLIDE_WIDTH - Inches(0.6), Inches(3.2))
    content_bg.fill.solid()
    content_bg.fill.fore_color.theme_color = get_theme_color(palette[0])
    content_bg.line.fill.background()

    row_width = (num_steps * box_width) + ((num_steps - 1) * Inches(0.6))
    start_x = MARGIN_LEFT + (CONTENT_WIDTH - row_width) / 2

    for i, step in enumerate(steps):
        left_offset = start_x + (box_width + Inches(0.6)) * i
        
        # Add Header Number Circle or Bar to denote step
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_offset, top_offset - Inches(0.4), box_width, Inches(0.5))
        hdr.fill.solid()
        hdr.fill.fore_color.theme_color = get_theme_color(palette[3])
        hdr.line.color.theme_color = get_theme_color(palette[3])
        hdr_tf = hdr.text_frame
        hdr_p = hdr_tf.paragraphs[0]
        hdr_p.text = str(i + 1).zfill(2)
        hdr_p.font.size = Pt(16)
        hdr_p.font.bold = True
        hdr_p.font.color.theme_color = get_theme_color(palette[0] if palette[0] != palette[3] else palette[1])
        hdr_p.alignment = PP_ALIGN.CENTER
        
        # Main Body Box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_offset, top_offset, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.theme_color = get_theme_color(palette[1])
        shape.line.color.theme_color = get_theme_color(palette[3])
        shape.line.width = Pt(1.0)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP # Keep top for infographics
        tf.margin_top = Inches(0.2)
        
        p_title = tf.paragraphs[0]
        p_title.text = step.get('title', f"Step {i+1}").upper()
        p_title.font.bold = True
        p_title.font.size = Pt(14)
        p_title.font.color.theme_color = get_theme_color(palette[2])
        p_title.alignment = PP_ALIGN.LEFT
        p_title.space_after = Pt(10)
        
        if 'desc' in step:
            p_desc = tf.add_paragraph()
            p_desc.text = step['desc']
            p_desc.font.size = Pt(12)
            p_desc.font.color.theme_color = get_theme_color(palette[2])
            p_desc.alignment = PP_ALIGN.LEFT
            
        # Add Connecting Arrow to next step
        if i < num_steps - 1:
            arrow_left = left_offset + box_width + Inches(0.1)
            arrow_top = top_offset + (box_height / 2) - Inches(0.2)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, Inches(0.4), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.theme_color = get_theme_color(palette[3])
            arrow.line.fill.background()

    prs.save(file_path)
    return "Infographic slide added."

@mcp.tool()
def add_chart_slide(file_path: str, title: str, chart_type: str, categories: list[str], series_data: list[dict], palette: list[str]) -> str:
    prs = Presentation(file_path)
    slide = prs.slides.add_slide(get_title_only_layout(prs))
    apply_title_styling(slide.shapes.title, title.upper(), palette)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series in series_data:
        chart_data.add_series(series['name'], series['values'])
        
    top = MARGIN_TOP
    left = MARGIN_LEFT + Inches(1.5)
    width = CONTENT_WIDTH - Inches(3)
    height = Inches(4.5)

    c_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
    if chart_type.upper() == 'BAR': c_type_enum = XL_CHART_TYPE.BAR_CLUSTERED
    if chart_type.upper() == 'LINE': c_type_enum = XL_CHART_TYPE.LINE
    if chart_type.upper() == 'PIE': c_type_enum = XL_CHART_TYPE.PIE

    slide.shapes.add_chart(c_type_enum, left, top, width, height, chart_data)
    prs.save(file_path)
    return f"Chart slide '{title}' added."

if __name__ == "__main__":
    try: mcp.run()
    except: pass

