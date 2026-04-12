"""
Stage 1 — Ingestion: Extract Design Tokens from a PowerPoint Slide Master.

Cracks open the user's master.pptx and extracts:
  - Corporate theme colors (HEX)
  - Font scheme (heading + body)
  - Slide dimensions
  - Layout catalog with placeholder positions
  - Composite-flattened background images per layout (Option B)

Outputs: design_tokens.json  +  assets/ folder with background PNGs.
"""

import os
import io
import json
import base64
import zipfile
import argparse
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu


# ─── XML Namespaces ───────────────────────────────────────────────────────────

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


# ─── Color Extraction ─────────────────────────────────────────────────────────

def extract_theme_colors(theme_xml: etree._Element) -> dict:
    """Extract all 12 theme color slots from clrScheme."""
    colors = {}
    elements = theme_xml.xpath("//a:clrScheme/*", namespaces=NSMAP)
    for el in elements:
        tag = el.tag.split("}")[-1]
        srgb = el.xpath(".//a:srgbClr/@val", namespaces=NSMAP)
        sys_clr = el.xpath(".//a:sysClr/@lastClr", namespaces=NSMAP)
        hex_val = srgb[0] if srgb else (sys_clr[0] if sys_clr else "000000")
        colors[tag] = f"#{hex_val}"
    return colors


# ─── Font Extraction ──────────────────────────────────────────────────────────

def extract_fonts(theme_xml: etree._Element) -> dict:
    """Extract major (heading) and minor (body) font faces."""
    major = theme_xml.xpath("//a:majorFont/a:latin/@typeface", namespaces=NSMAP)
    minor = theme_xml.xpath("//a:minorFont/a:latin/@typeface", namespaces=NSMAP)
    return {
        "heading": major[0] if major else "Calibri",
        "body": minor[0] if minor else "Calibri",
    }


# ─── Layout Catalog ───────────────────────────────────────────────────────────

def _normalize_ph_type(ph_type_str: str) -> str:
    """Normalize placeholder type enum to a clean lowercase name."""
    # ph_type_str is like 'TITLE (15)' or 'BODY (2)'
    name = str(ph_type_str).split('(')[0].strip().lower()
    name = name.replace(' ', '_')
    # Map common aliases
    aliases = {
        'center_title': 'title',
        'subtitle': 'subtitle',
        'vertical_body': 'body',
        'vertical_title': 'title',
        'vertical_object': 'body',
        'object': 'body',
    }
    return aliases.get(name, name)


def extract_layouts(prs: Presentation) -> list:
    """Extract layout names, indices, placeholder geometry, and type signatures."""
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        signature = {}  # Count of each placeholder type
        for ph in layout.placeholders:
            ph_type_raw = str(ph.placeholder_format.type)
            ph_type_clean = _normalize_ph_type(ph_type_raw)
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "type": ph_type_raw,
                "type_clean": ph_type_clean,
                "name": ph.name,
                "x": round(ph.left / Emu(914400), 2),   # EMU → inches
                "y": round(ph.top / Emu(914400), 2),
                "w": round(ph.width / Emu(914400), 2),
                "h": round(ph.height / Emu(914400), 2),
            })
            signature[ph_type_clean] = signature.get(ph_type_clean, 0) + 1

        layouts.append({
            "index": i,
            "name": layout.name,
            "placeholders": placeholders,
            "signature": signature,
        })
    return layouts


# ─── Background Composite Extraction (Option B) ───────────────────────────────

def _get_all_media(pptx_path: str) -> dict:
    """Read all media files from the pptx ZIP into memory as {name: bytes}."""
    media = {}
    with zipfile.ZipFile(pptx_path, "r") as z:
        for entry in z.namelist():
            if "ppt/media/" in entry:
                media[entry] = z.read(entry)
    return media


def _resolve_rId(part, rId: str):
    """Resolve a relationship ID to a target part's blob."""
    try:
        related = part.related_parts.get(rId)
        if related:
            return related.blob
    except Exception:
        pass
    return None


def _render_shape_to_image(shape_xml, part, slide_w_px, slide_h_px, media_cache):
    """
    Attempt to render a shape element (image fill, picture, solid fill) into a
    PIL Image layer of size (slide_w_px, slide_h_px).
    Returns an RGBA Image or None.
    """
    layer = None

    # Position and size
    off = shape_xml.find(".//p:spPr/a:xfrm/a:off", NSMAP)
    if off is None:
        off = shape_xml.find(".//a:xfrm/a:off", NSMAP)
    ext = shape_xml.find(".//p:spPr/a:xfrm/a:ext", NSMAP)
    if ext is None:
        ext = shape_xml.find(".//a:xfrm/a:ext", NSMAP)

    if off is not None and ext is not None:
        x = int(off.get("x", "0"))
        y = int(off.get("y", "0"))
        cx = int(ext.get("cx", "0"))
        cy = int(ext.get("cy", "0"))
    else:
        x, y, cx, cy = 0, 0, 0, 0

    # Convert EMU to pixels (at 96 DPI, 1 inch = 914400 EMU)
    ppi = 96
    emu_per_px = 914400 / ppi
    px_x = int(x / emu_per_px)
    px_y = int(y / emu_per_px)
    px_w = int(cx / emu_per_px) if cx > 0 else slide_w_px
    px_h = int(cy / emu_per_px) if cy > 0 else slide_h_px

    # Case 1: Image fill (blipFill)
    blip = shape_xml.find(".//a:blipFill/a:blip", NSMAP)
    if blip is None:
        blip = shape_xml.find(".//p:blipFill/a:blip", NSMAP)

    if blip is not None:
        r_embed = blip.get(f"{{{NSMAP['r']}}}embed")
        if r_embed:
            img_blob = _resolve_rId(part, r_embed)
            if img_blob:
                try:
                    img = Image.open(io.BytesIO(img_blob)).convert("RGBA")
                    img = img.resize((px_w if px_w > 0 else slide_w_px,
                                      px_h if px_h > 0 else slide_h_px),
                                     Image.LANCZOS)
                    layer = Image.new("RGBA", (slide_w_px, slide_h_px), (0, 0, 0, 0))
                    layer.paste(img, (px_x, px_y))
                    return layer
                except Exception:
                    pass

    # Case 2: Solid fill on the shape
    solid = shape_xml.find(".//p:spPr/a:solidFill/a:srgbClr", NSMAP)
    if solid is not None:
        hex_val = solid.get("val", "000000")
        r, g, b = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
        layer = Image.new("RGBA", (slide_w_px, slide_h_px), (r, g, b, 255))
        return layer

    return None


def _composite_layout_background(prs, layout, pptx_path, slide_w_px, slide_h_px):
    """
    Composite all background layers of a slide layout + its master into a
    single flattened RGBA image.
    """
    canvas = Image.new("RGBA", (slide_w_px, slide_h_px), (255, 255, 255, 255))

    # We work directly with the XML to capture all shapes, including non-placeholder ones
    # which python-pptx hides (background rectangles, decorative images, etc.)

    # Layer 1: Slide Master shapes (bottom layer)
    master = layout.slide_master
    master_part = master.part
    master_xml = master_part._element

    for sp in master_xml.iter(f"{{{NSMAP['p']}}}sp"):
        layer = _render_shape(sp, master_part, slide_w_px, slide_h_px)
        if layer:
            canvas = Image.alpha_composite(canvas, layer)

    for pic in master_xml.iter(f"{{{NSMAP['p']}}}pic"):
        layer = _render_shape(pic, master_part, slide_w_px, slide_h_px)
        if layer:
            canvas = Image.alpha_composite(canvas, layer)

    # Layer 2: Slide Layout shapes (top layer, overrides master)
    layout_part = layout.part
    layout_xml = layout_part._element

    for sp in layout_xml.iter(f"{{{NSMAP['p']}}}sp"):
        layer = _render_shape(sp, layout_part, slide_w_px, slide_h_px)
        if layer:
            canvas = Image.alpha_composite(canvas, layer)

    for pic in layout_xml.iter(f"{{{NSMAP['p']}}}pic"):
        layer = _render_shape(pic, layout_part, slide_w_px, slide_h_px)
        if layer:
            canvas = Image.alpha_composite(canvas, layer)

    return canvas


def _render_shape(shape_xml, part, slide_w_px, slide_h_px):
    """
    Universal shape renderer: handles pictures and shapes with image/solid fills.
    Returns an RGBA Image layer of size (slide_w_px, slide_h_px) or None.
    """
    layer = None
    ppi = 96
    emu_per_px = 914400 / ppi

    # Position and size from xfrm
    xfrm = shape_xml.find(".//p:spPr/a:xfrm", NSMAP)
    if xfrm is None:
        xfrm = shape_xml.find(".//a:xfrm", NSMAP)
    if xfrm is None:
        # Try under the nvSpPr sibling path — it might be directly under the shape
        xfrm = shape_xml.find("p:spPr/a:xfrm", NSMAP)

    if xfrm is not None:
        off = xfrm.find("a:off", NSMAP)
        ext = xfrm.find("a:ext", NSMAP)
        if off is not None and ext is not None:
            x = int(off.get("x", "0"))
            y = int(off.get("y", "0"))
            cx = int(ext.get("cx", "0"))
            cy = int(ext.get("cy", "0"))
        else:
            return None
    else:
        return None

    px_x = int(x / emu_per_px)
    px_y = int(y / emu_per_px)
    px_w = int(cx / emu_per_px) if cx > 0 else slide_w_px
    px_h = int(cy / emu_per_px) if cy > 0 else slide_h_px

    if px_w <= 0 or px_h <= 0:
        return None

    # Case 1: blipFill (image fill — used by pictures and shapes with image fills)
    blip = shape_xml.find(".//a:blipFill/a:blip", NSMAP)
    if blip is None:
        blip = shape_xml.find(".//p:blipFill/a:blip", NSMAP)
    if blip is None:
        blip = shape_xml.find(".//a:blip", NSMAP)

    if blip is not None:
        r_embed = blip.get(f"{{{NSMAP['r']}}}embed")
        if r_embed:
            try:
                related_part = part.related_parts.get(r_embed)
                if related_part is None:
                    related_part = part.rels[r_embed].target_part
                img_blob = related_part.blob
                img = Image.open(io.BytesIO(img_blob)).convert("RGBA")
                img = img.resize((px_w, px_h), Image.LANCZOS)
                layer = Image.new("RGBA", (slide_w_px, slide_h_px), (0, 0, 0, 0))
                layer.paste(img, (px_x, px_y))
                return layer
            except Exception:
                pass

    # Case 2: Solid fill
    solid = shape_xml.find(".//p:spPr/a:solidFill/a:srgbClr", NSMAP)
    if solid is None:
        solid = shape_xml.find("p:spPr/a:solidFill/a:srgbClr", NSMAP)
    if solid is not None:
        hex_val = solid.get("val", "000000")
        # Check for alpha
        alpha_el = solid.find("a:alpha", NSMAP)
        alpha = 255
        if alpha_el is not None:
            alpha = int(int(alpha_el.get("val", "100000")) / 1000 * 255 / 100)
        r_c, g_c, b_c = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
        fill_img = Image.new("RGBA", (px_w, px_h), (r_c, g_c, b_c, alpha))
        layer = Image.new("RGBA", (slide_w_px, slide_h_px), (0, 0, 0, 0))
        layer.paste(fill_img, (px_x, px_y))
        return layer

    # Case 3: Gradient fill — render as linear gradient from first to last stop
    grad = shape_xml.find(".//p:spPr/a:gradFill", NSMAP)
    if grad is None:
        grad = shape_xml.find("p:spPr/a:gradFill", NSMAP)
    if grad is not None:
        stops = grad.findall(".//a:gs/a:srgbClr", NSMAP)
        if len(stops) >= 2:
            c1_hex = stops[0].get("val", "000000")
            c2_hex = stops[-1].get("val", "FFFFFF")
            r1, g1, b1 = int(c1_hex[0:2], 16), int(c1_hex[2:4], 16), int(c1_hex[4:6], 16)
            r2, g2, b2 = int(c2_hex[0:2], 16), int(c2_hex[2:4], 16), int(c2_hex[4:6], 16)
            fill_img = Image.new("RGBA", (px_w, px_h))
            for row in range(px_h):
                ratio = row / max(px_h - 1, 1)
                r_v = int(r1 + (r2 - r1) * ratio)
                g_v = int(g1 + (g2 - g1) * ratio)
                b_v = int(b1 + (b2 - b1) * ratio)
                for col in range(px_w):
                    fill_img.putpixel((col, row), (r_v, g_v, b_v, 255))
            layer = Image.new("RGBA", (slide_w_px, slide_h_px), (0, 0, 0, 0))
            layer.paste(fill_img, (px_x, px_y))
            return layer

    return None


# ─── Main Pipeline ─────────────────────────────────────────────────────────────

def ingest(template_path: str, output_dir: str = "output") -> str:
    """
    Full ingestion pipeline. Returns the path to design_tokens.json.
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    # Create output directories
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    prs = Presentation(template_path)

    # Dimensions
    slide_w_in = round(prs.slide_width / 914400, 2)
    slide_h_in = round(prs.slide_height / 914400, 2)
    ppi = 96
    slide_w_px = int(slide_w_in * ppi)
    slide_h_px = int(slide_h_in * ppi)

    # Theme XML (from first slide master)
    sm = prs.slide_masters[0]
    theme_part = sm.part.part_related_by(RT.THEME)
    theme_xml = etree.fromstring(theme_part.blob)

    # Extract tokens
    colors = extract_theme_colors(theme_xml)
    fonts = extract_fonts(theme_xml)
    layouts = extract_layouts(prs)

    # Extract composite backgrounds per layout
    backgrounds = {}
    print(f"Extracting {len(prs.slide_layouts)} layout backgrounds...")
    for i, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name.lower().replace(" ", "_").replace("__", "_")
        print(f"  Compositing layout {i}: {layout.name}...")
        try:
            bg_img = _composite_layout_background(prs, layout, template_path, slide_w_px, slide_h_px)
            bg_filename = f"bg_{layout_name}.png"
            bg_path = os.path.join(assets_dir, bg_filename)
            bg_img.convert("RGB").save(bg_path, "PNG", optimize=True)
            backgrounds[layout.name] = os.path.join("assets", bg_filename).replace("\\", "/")
            print(f"    -> Saved {bg_filename}")
        except Exception as e:
            print(f"    -> WARNING: Could not composite background for '{layout.name}': {e}")
            backgrounds[layout.name] = None

    # Assemble design tokens
    tokens = {
        "template_name": os.path.basename(template_path),
        "dimensions": {
            "width": slide_w_in,
            "height": slide_h_in,
            "width_px": slide_w_px,
            "height_px": slide_h_px,
        },
        "colors": colors,
        "fonts": fonts,
        "layouts": layouts,
        "backgrounds": backgrounds,
    }

    tokens_path = os.path.join(output_dir, "design_tokens.json")
    with open(tokens_path, "w", encoding="utf-8") as f:
        json.dump(tokens, f, indent=2, ensure_ascii=False)

    print(f"\nDesign tokens saved to: {tokens_path}")
    print(f"  Colors: {len(colors)} slots")
    print(f"  Fonts: heading='{fonts['heading']}', body='{fonts['body']}'")
    print(f"  Layouts: {len(layouts)}")
    print(f"  Backgrounds: {sum(1 for v in backgrounds.values() if v)} composited")
    return tokens_path


# ─── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Stage 1 — Ingest: Extract design tokens from a PowerPoint Slide Master."
    )
    parser.add_argument("--template", required=True, help="Path to the master .pptx template")
    parser.add_argument("--output", default="output", help="Output directory (default: output/)")
    args = parser.parse_args()

    ingest(args.template, args.output)
