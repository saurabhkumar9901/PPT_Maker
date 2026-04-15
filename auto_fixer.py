"""
Stage 4 -- Auto-Fixer: Programmatic Post-Generation Linter
Applies strict mathematical rules to fix common mistakes post-rendering.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

def run_fixes(pptx_path: str, tokens: dict = None, output_path: str = None) -> str:
    """Run programmatic auto-fixes over a generated presentation."""
    if not os.path.exists(pptx_path):
        return pptx_path

    out_path = output_path or pptx_path
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"AutoFixer could not read file: {e}")
        return pptx_path

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for i, slide in enumerate(prs.slides):
        content_shapes = []
        for shape in slide.shapes:
            # Gather valid content shapes for later layout spacing calculation
            if shape.has_text_frame or shape.shape_type == 14: # 14 is placeholder
                content_shapes.append(shape)
                
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            has_solid_fill = False
            
            # Rule 1: Zero-Margin Enforcement on transparent boxes
            try:
                # Type 1 is SOLID, 5 is BACKGROUND
                if shape.fill.type is not None and getattr(shape.fill.type, 'name', '') == 'SOLID':
                    has_solid_fill = True
            except Exception:
                pass

            if not has_solid_fill:
                tf.margin_left = 0
                tf.margin_right = 0
                tf.margin_top = 0
                tf.margin_bottom = 0

            # Rule 2: Table/Matrix Middle Alignment
            # Heuristically, matrix/table boxes are typically short and have short text.
            is_probably_cell = False
            if shape.height and shape.width:
                if shape.height < Inches(1.0) and shape.width < Inches(3.0) and has_solid_fill:
                    is_probably_cell = True

            if is_probably_cell:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.CENTER

            # Rule 3: Bounding Box Overflow Failsafe Shrinker
            if shape.left is not None and shape.width is not None:
                if shape.left + shape.width > slide_w + Inches(0.2):
                    # It's spilling off the right! Pull font sizes down actively.
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                run.font.size = Pt(max(6, run.font.size.pt * 0.85)) # 15% reduction

        # Rule 4: Negative Space Distributer
        # If the slide only has a title and a tiny bit of content, bump it down slightly
        if len(content_shapes) in [2, 3] and slide_h > Inches(5):
            lowest_y = 0
            for s in content_shapes:
                if s.top and s.height:
                    bottom = s.top + s.height
                    if bottom > lowest_y:
                        lowest_y = bottom
            
            # If the lowest element doesn't even reach 50% down the slide, bump the body shapes down
            if lowest_y > 0 and lowest_y < slide_h * 0.5:
                # We skip the title shape (usually high up and wide)
                for s in content_shapes:
                    if s.top and s.top > Inches(1.5): # Probably a body shape
                        s.top = int(s.top + Inches(1.0))

    try:
        prs.save(out_path)
        print(f"Auto-fixed presentation saved to {out_path}")
        return out_path
    except Exception as e:
        print(f"Error saving auto-fixed presentation: {e}")
        return pptx_path

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        run_fixes(sys.argv[1])
