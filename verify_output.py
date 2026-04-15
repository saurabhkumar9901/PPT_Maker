from pptx import Presentation
import sys

prs = Presentation(sys.argv[1])
print("==== SLIDE 6 ====")
slide6 = prs.slides[5]
for s in slide6.shapes:
    text = ""
    font_color = ""
    if s.has_text_frame:
        text = " ".join([p.text for p in s.text_frame.paragraphs])
        for p in s.text_frame.paragraphs:
            if p.runs and p.runs[0].font.color.type is not None:
                font_color = str(p.runs[0].font.color.rgb) if hasattr(p.runs[0].font.color, 'rgb') else 'theme/default'
    print(f"Type: {s.shape_type}, Name: {s.name}, Text: {text[:60]}, Color: {font_color}")

print("==== SLIDE 12 ====")
slide12 = prs.slides[11]
for s in slide12.shapes[-6:]:
    text = ""
    font_color = ""
    if s.has_text_frame:
        text = " ".join([p.text for p in s.text_frame.paragraphs])
        for p in s.text_frame.paragraphs:
            for r in p.runs:
               if r.font.color.type is not None and hasattr(r.font.color, 'rgb'):
                   font_color = str(r.font.color.rgb)
    print(f"Type: {s.shape_type}, Name: {s.name}, W: {s.width}, H: {s.height}, Text: {text[:60]}, Color: {font_color}")
