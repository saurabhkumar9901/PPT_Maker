from pptx import Presentation
import sys

def dump_slide(prs_path, slide_idx=1): # Slide 2 usually content
    print(f"\n--- Analyzing Slide {slide_idx+1} of {prs_path} ---")
    try:
        prs = Presentation(prs_path)
    except Exception as e:
        print(f"Could not load {prs_path}: {e}")
        return
        
    if slide_idx >= len(prs.slides):
        print("Slide index out of range")
        return
        
    slide = prs.slides[slide_idx]
    print(f"Total shapes: {len(slide.shapes)}")
    
    for i, shape in enumerate(slide.shapes):
        print(f"\nShape {i}: {shape.shape_type} | Name: {shape.name}")
        try:
            print(f"  Pos/Size: L={shape.left.inches:.2f}, T={shape.top.inches:.2f}, W={shape.width.inches:.2f}, H={shape.height.inches:.2f}")
        except: pass
        
        if shape.has_text_frame:
            tf = shape.text_frame
            print(f"  Text: {tf.text[:50]}...")
            if tf.paragraphs:
                p = tf.paragraphs[0]
                if p.runs:
                    font = p.runs[0].font
                    size = font.size.pt if font.size else 'Default'
                    color = font.color.rgb if font.color and hasattr(font.color, 'rgb') else 'Default'
                    bold = font.bold
                    print(f"  Font Formatting: size={size}, color={color}, bold={bold}")
        elif shape.has_chart:
            print(f"  Chart Type: {shape.chart.chart_type}")
        elif shape.has_table:
            print(f"  Table: {len(shape.table.rows)} rows x {len(shape.table.columns)} cols")

if __name__ == "__main__":
    sample = r"d:\PPT Maker\Sample Files\Sample Files\Accenture Tech Acquisition Analysis\Accenture Tech Acquisition Analysis.pptx"
    output = r"d:\PPT Maker\Output_Accenture_Deck.pptx"
    
    dump_slide(sample, 0) # Cover
    dump_slide(sample, 1) # First content
    dump_slide(sample, 2) # Second content
    
    dump_slide(output, 0)
    dump_slide(output, 1)
    dump_slide(output, 2)
