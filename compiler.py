"""
Stage 3 -- Compiler (Python): Design Tokens + UI Plan -> .pptx

Uses python-pptx to create slides from the NATIVE template (preserving
all background images, EMF vectors, SVG decorations, gradients, etc.),
then dynamically places content elements using precise coordinate math.

Design principles applied (from Common Mistakes guide):
  - Section labels on every content slide
  - Footer with report name on every content slide
  - Zero internal margins on transparent text boxes
  - Separator lines between cards/columns for visual rhythm
  - Numbered badges on grid items
  - Custom shape-based tables (not native table objects)
  - Proper font sizing: 30pt numbers, 9pt labels
  - Consistent spacing throughout

Usage: python compiler.py <tokens_path> <plan_path> <template_path> <output_path>
"""

import os
import sys
import json
import math
import argparse
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ─── Configuration ────────────────────────────────────────────────────────────

MARGIN_L = Inches(0.38)
MARGIN_R = Inches(0.38)
TITLE_Y  = Inches(0.60)
TITLE_H  = Inches(0.85)
SECTION_Y = Inches(0.22)
SECTION_H = Inches(0.28)
CONTENT_TOP = Inches(1.60)
FOOTER_Y = Inches(7.23)
FOOTER_H = Inches(0.20)
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.50)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
CONTENT_H = SLIDE_H - CONTENT_TOP - Inches(0.50)  # leave room for footer


# ─── Helpers ──────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_str):
    h = (hex_str or '#333333').lstrip('#')
    if len(h) < 6: h = h.ljust(6, '0')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _zero_margins(tf):
    """Set all internal margins on a text frame to zero."""
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

def _small_margins(tf, lr=Inches(0.08), tb=Inches(0.04)):
    tf.margin_left = lr
    tf.margin_right = lr
    tf.margin_top = tb
    tf.margin_bottom = tb

def get_layout(prs, layout_type):
    name_map = {
        'cover': ['cover', '1_cover', '2_cover', '0_title company'],
        'divider': ['divider', 'c_section blue', 'section'],
        'content': ['title only', '1_e_title, subtitle and body', 'blank'],
        'chart': ['title only', '1_e_title, subtitle and body', 'blank'],
        'thank_you': ['1_thank you', 'thank you', 'thank_you'],
    }
    candidates = name_map.get(layout_type, name_map['content'])
    for cand in candidates:
        for layout in prs.slide_layouts:
            if cand in layout.name.lower():
                return layout
    return prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None


# ─── Slide Furniture ──────────────────────────────────────────────────────────

def add_section_label(slide, section_text, tokens):
    """Small accent-colored section label at top-left (e.g. 'SECTION 1')."""
    w = Inches(1.50)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, SECTION_Y, w, SECTION_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(tokens['colors']['dk1'])
    shape.line.fill.background()
    # Tight radius
    tf = shape.text_frame
    _zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = section_text.upper()
    p.font.name = tokens['fonts']['heading']
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(tokens['colors']['lt1'])
    p.alignment = PP_ALIGN.CENTER

def add_footer(slide, tokens, page_num, report_name=None):
    """Bottom footer with report name and page number."""
    name = report_name or tokens.get('template_name', '')
    # Clean label
    name = name.replace('.pptx', '').replace('.ppt', '').replace('Template_', '').replace('template_', '').replace('_', ' ')
    # Report name - left
    txBox = slide.shapes.add_textbox(MARGIN_L, FOOTER_Y, Inches(9.0), FOOTER_H)
    tf = txBox.text_frame
    _zero_margins(tf)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = name
    p.font.name = tokens['fonts']['body']
    p.font.size = Pt(7)
    p.font.color.rgb = hex_to_rgb(tokens['colors'].get('dk2', '#4A6080'))
    p.alignment = PP_ALIGN.LEFT
    # Page number - right
    txBox2 = slide.shapes.add_textbox(Inches(9.94), FOOTER_Y, Inches(3.0), FOOTER_H)
    tf2 = txBox2.text_frame
    _zero_margins(tf2)
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num)
    p2.font.name = tokens['fonts']['body']
    p2.font.size = Pt(7)
    p2.font.color.rgb = hex_to_rgb(tokens['colors'].get('dk2', '#4A6080'))
    p2.alignment = PP_ALIGN.RIGHT


def add_title(slide, title_text, tokens):
    if not title_text:
        return
    # Use existing title placeholder if available
    if slide.shapes.title is not None:
        shape = slide.shapes.title
        shape.text = title_text
        # Force alignment to avoid overlapping with section badges
        shape.left = int(MARGIN_L)
        shape.top = int(SECTION_Y + SECTION_H + Inches(0.20))
        shape.width = int(CONTENT_W)
        shape.height = int(TITLE_H)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for p in tf.paragraphs:
            p.font.name = tokens['fonts']['heading']
            p.font.size = Pt(22) if len(title_text) <= 50 else Pt(18)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])
            p.alignment = PP_ALIGN.LEFT
        return
    # Fallback text box
    txBox = slide.shapes.add_textbox(MARGIN_L, TITLE_Y, CONTENT_W, TITLE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _zero_margins(tf)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = tokens['fonts']['heading']
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])
    p.alignment = PP_ALIGN.LEFT


# ─── Primitive helpers ────────────────────────────────────────────────────────

from pptx.oxml.xmlchemy import OxmlElement

def add_text_box(slide, x, y, w, h, text, font_name, font_size, color,
                 bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, wrap=True, zero_margin=True):
    txBox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    if zero_margin:
        _zero_margins(tf)
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_card(slide, x, y, w, h, fill_color, border_color=None, border_w=Pt(1), alpha=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if alpha is not None:
        a = OxmlElement('a:alpha')
        a.set('val', str(alpha))
        shape.fill.fore_color._color._xClr.append(a)
        
    if shadow:
        try:
            from pptx.oxml import parse_xml
            shadow_xml = """<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                <a:outerShdw blurRad="100000" dist="40000" dir="5400000" algn="b" rotWithShape="0">
                    <a:srgbClr val="000000">
                        <a:alpha val="15000"/>
                    </a:srgbClr>
                </a:outerShdw>
            </a:effectLst>"""
            effect_lst = parse_xml(shadow_xml)
            shape.element.spPr.append(effect_lst)
        except Exception:
            pass

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape

def add_hline(slide, x, y, w, color, thickness=Pt(1)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_vline(slide, x, y, h, color, thickness=Pt(1)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(thickness), int(h))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_numbered_badge(slide, x, y, num, tokens, size=Inches(0.40)):
    """Circle badge with number, accent colored."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(size), int(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(tokens['colors']['accent1'])
    shape.line.fill.background()
    tf = shape.text_frame
    _zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"{num:02d}" if isinstance(num, int) else str(num)
    p.font.name = tokens['fonts']['heading']
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(tokens['colors']['lt1'])
    p.alignment = PP_ALIGN.CENTER

def _add_bullet_list(slide, x, y, w, h, items, tokens, font_size=Pt(11)):
    """Add a bulleted text box with styled items."""
    txBox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _zero_margins(tf)

    for i, item in enumerate(items):
        if isinstance(item, str):
            text = item
            bold_prefix = None
        else:
            text = item.get('text', '')
            bold_prefix = item.get('bold_prefix', None)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)

        if bold_prefix:
            run = p.add_run()
            run.text = bold_prefix + '  '
            run.font.name = tokens['fonts']['heading']
            run.font.size = font_size
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])

        run = p.add_run()
        run.text = text
        run.font.name = tokens['fonts']['body']
        run.font.size = font_size
        run.font.color.rgb = hex_to_rgb(tokens['colors']['dk2'])

        # Bullet via XML
        pPr = p._pPr
        if pPr is None:
            p._p.insert(0, p._p.makeelement(f'{{{A_NS}}}pPr', {}))
            pPr = p._pPr
        buChar = etree.SubElement(pPr, f'{{{A_NS}}}buChar')
        buChar.set('char', '\u2022')

    return txBox


# ─── Element Renderers ────────────────────────────────────────────────────────

def render_grid(slide, element, tokens):
    """Infographic-style grid with distinct header strips and centered alignment."""
    items = element.get('items', [])
    cols = min(element.get('columns', 3), 4)
    if not items: return

    gap = Inches(0.20)
    rows = -(-len(items) // cols)
    row_gap = Inches(0.20)
    card_w = int((CONTENT_W - gap * (cols - 1)) / cols)
    card_h = int((CONTENT_H - row_gap * (rows - 1)) / rows)
    card_h = min(card_h, int(CONTENT_H * 0.8))

    badge_size = Inches(0.35)
    header_strip_h = Inches(0.60)

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        x = int(MARGIN_L) + col * (card_w + int(gap))
        y = int(CONTENT_TOP) + row * (card_h + int(row_gap))

        # Main Card Body with subtle accent tint (instead of faded lt2)
        accent = tokens['colors']['accent1']
        add_card(slide, x, y, card_w, card_h,
                hex_to_rgb(accent), 
                border_color=hex_to_rgb(accent), border_w=Pt(1.5), alpha=15000, shadow=True)

        # Header Strip (solid color)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, int(header_strip_h))
        strip.fill.solid()
        strip.fill.fore_color.rgb = hex_to_rgb(accent)
        strip.line.fill.background()

        # Numbered badge (offset overlapping strip)
        add_numbered_badge(slide, x + Inches(0.15), y + Inches(0.12), i + 1, tokens, badge_size)

        # Heading (inside the header strip)
        add_text_box(slide, x + Inches(0.15) + badge_size + Inches(0.08), y + Inches(0.12),
                    card_w - Inches(0.38) - badge_size, badge_size,
                    item.get('heading', ''), tokens['fonts']['heading'], Pt(11),
                    hex_to_rgb(tokens['colors']['lt1']), bold=True,
                    valign=MSO_ANCHOR.MIDDLE)

        # Body text (middle-aligned in the remaining space for better distribution)
        add_text_box(slide, x + Inches(0.15), y + int(header_strip_h) + Inches(0.10),
                    card_w - Inches(0.30), card_h - int(header_strip_h) - Inches(0.20),
                    item.get('body', ''), tokens['fonts']['body'], Pt(10),
                    hex_to_rgb(tokens['colors']['dk2']), 
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def render_timeline(slide, element, tokens):
    """High-contrast interlocking timeline process map."""
    steps = element.get('steps', [])
    if not steps: return
    n = len(steps)
    gap = Inches(0.10)
    connector_w = Inches(0.20)
    box_w = int((CONTENT_W - connector_w * (n - 1) - gap * (n - 1) * 2) / n)
    box_w = min(box_w, Inches(2.8))
    box_h = int(CONTENT_H * 0.65)
    badge_h = Inches(0.40)
    total_w = n * box_w + (n - 1) * (int(gap) * 2 + int(connector_w))
    start_x = int(MARGIN_L) + int((CONTENT_W - total_w) / 2)
    mid_y = int(CONTENT_TOP) + Inches(0.4)

    for i, step in enumerate(steps):
        x = start_x + i * (box_w + int(gap) * 2 + int(connector_w))

        # Timeline track connector logic (draw line under badges)
        if i < n - 1:
            ax = x + box_w
            ay = mid_y + int(badge_h / 2) - Inches(0.05)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, int(connector_w) + int(gap)*2, Inches(0.10))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = hex_to_rgb(tokens['colors']['accent1'])
            arrow.line.fill.background()

        # Step Label Badge (solid accent color as Pill Shape)
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, mid_y, box_w, int(badge_h))
        badge.fill.solid()
        badge.fill.fore_color.rgb = hex_to_rgb(tokens['colors']['accent1'])
        badge.line.fill.background()
        
        try:
            from pptx.oxml import parse_xml
            adj_xml = '<a:avLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:gd name="adj" fmla="val 50000"/></a:avLst>'
            prstGeom = badge.element.spPr.prstGeom
            if prstGeom.avLst is not None:
                prstGeom.remove(prstGeom.avLst)
            prstGeom.append(parse_xml(adj_xml))
        except Exception:
            pass
            
        tf = badge.text_frame
        _zero_margins(tf)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].text = step.get('label', f'{i+1:02d}').upper()
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = hex_to_rgb(tokens['colors']['lt1'])
        tf.paragraphs[0].font.name = tokens['fonts']['heading']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Box (Content area below badge)
        box_y = mid_y + int(badge_h) + Inches(0.10)
        add_card(slide, x, box_y, box_w, int(box_h),
                hex_to_rgb(tokens['colors']['lt2']),
                border_color=hex_to_rgb(tokens['colors']['accent1']), border_w=Pt(1.5), shadow=True)

        # Title
        add_text_box(slide, x + Inches(0.12), box_y + Inches(0.10), box_w - Inches(0.24), Inches(0.35),
                    step.get('title', '').upper(), tokens['fonts']['heading'], Pt(11),
                    hex_to_rgb(tokens['colors']['dk1']), bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Subtle separator
        add_hline(slide, x + int((box_w - Inches(1.0))/2), box_y + Inches(0.50),
                 Inches(1.0), hex_to_rgb(tokens['colors']['accent1']), Pt(2))

        # Description
        add_text_box(slide, x + Inches(0.15), box_y + Inches(0.60), box_w - Inches(0.30), int(box_h) - Inches(0.70),
                    step.get('description', ''), tokens['fonts']['body'], Pt(10),
                    hex_to_rgb(tokens['colors']['dk2']), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)


def render_hero(slide, element, tokens):
    query = element.get('image_query')
    if query:
        add_image_placeholder(slide, MARGIN_L, CONTENT_TOP + Inches(0.2), CONTENT_W, CONTENT_H - Inches(0.4), query, tokens, show_icon_text=False)
        
    # Draw content on a slightly darker opaque card so text is readable over "image"
    if query:
        add_card(slide, MARGIN_L + Inches(1.0), CONTENT_TOP + Inches(1.0), 
                 CONTENT_W - Inches(2.0), CONTENT_H - Inches(1.0), 
                 hex_to_rgb(tokens['colors']['lt1']), alpha=90000, border_color=None)
    
    add_text_box(slide, MARGIN_L + Inches(0.5), CONTENT_TOP + Inches(1.5 if query else 0.6),
                CONTENT_W - Inches(1.0), Inches(1.2),
                element.get('heading', ''), tokens['fonts']['heading'], Pt(32),
                hex_to_rgb(tokens['colors']['dk1']), bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    line_w = Inches(2.0)
    add_hline(slide, int(MARGIN_L) + int((CONTENT_W - line_w) / 2),
             int(CONTENT_TOP) + Inches(2.7 if query else 2.0), line_w,
             hex_to_rgb(tokens['colors']['dk1']), Pt(2))

    add_text_box(slide, MARGIN_L + Inches(1.0), CONTENT_TOP + Inches(2.95 if query else 2.25),
                CONTENT_W - Inches(2.0), CONTENT_H - Inches(3.2 if query else 2.5),
                element.get('body', ''), tokens['fonts']['body'], Pt(12),
                hex_to_rgb(tokens['colors']['dk2']),
                align=PP_ALIGN.CENTER)


def render_bullets(slide, element, tokens):
    items = element.get('items', [])
    if not items: return
    _add_bullet_list(slide, MARGIN_L + Inches(0.10), CONTENT_TOP + Inches(0.15),
                     CONTENT_W - Inches(0.20), CONTENT_H - Inches(0.30),
                     items, tokens, Pt(11))


def render_chart(slide, element, tokens):
    chart_type_map = {
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE,
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'area': XL_CHART_TYPE.AREA,
    }
    ct = element.get('chart_type', 'bar').lower()
    chart_type = chart_type_map.get(ct, XL_CHART_TYPE.BAR_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = element.get('categories', [])
    for series in element.get('series', []):
        chart_data.add_series(series.get('name', ''), series.get('values', []))

    chart_frame = slide.shapes.add_chart(
        chart_type,
        MARGIN_L + Inches(0.25), CONTENT_TOP + Inches(0.10),
        CONTENT_W - Inches(0.50), CONTENT_H - Inches(0.20),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    try:
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.color.rgb = hex_to_rgb(tokens['colors']['dk2'])
    except Exception:
        pass

    try:
        for series in chart.series:
            series.has_data_labels = True
            try:
                series.data_labels.font.color.rgb = hex_to_rgb(tokens['colors']['dk2'])
            except Exception:
                pass
    except Exception:
        pass


def render_table(slide, element, tokens):
    """Custom shape-based table with line separators (not native table)."""
    headers = element.get('headers', [])
    rows_data = element.get('rows', [])
    if not headers: return

    n_cols = len(headers)
    n_rows = len(rows_data)
    col_w = int(CONTENT_W / n_cols)
    row_h = Inches(0.55)
    header_h = Inches(0.40)
    x_start = int(MARGIN_L)
    y_start = int(CONTENT_TOP) + Inches(0.10)

    # Header row
    for j, h in enumerate(headers):
        x = x_start + j * col_w
        add_text_box(slide, x, y_start, col_w, header_h,
                    h, tokens['fonts']['heading'], Pt(12),
                    hex_to_rgb(tokens['colors']['dk1']), bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Header underline
    add_hline(slide, x_start, y_start + header_h,
             int(CONTENT_W), hex_to_rgb(tokens['colors']['accent1']), Pt(2))

    # Data rows
    for ri, row in enumerate(rows_data):
        cells = row.get('cells', []) if isinstance(row, dict) else row
        y = y_start + header_h + Inches(0.08) + ri * row_h

        # Row data cells
        for j, cell_text in enumerate(cells):
            if j >= n_cols: break
            x = x_start + j * col_w
            is_first_col = j == 0
            add_text_box(slide, x, y, col_w, row_h,
                        str(cell_text), tokens['fonts']['body'] if not is_first_col else tokens['fonts']['heading'],
                        Pt(11), hex_to_rgb(tokens['colors']['dk2']),
                        bold=is_first_col, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Row separator line
        if ri < n_rows - 1:
            line_y = y + row_h
            add_hline(slide, x_start + Inches(0.15), line_y,
                     int(CONTENT_W) - Inches(0.30),
                     hex_to_rgb(tokens['colors'].get('dk2', '#CCCCCC')), Pt(0.5))

    # Column separator badges
    for j in range(1, n_cols):
        cx = x_start + j * col_w - Inches(0.12)
        for ri in range(n_rows):
            cy = y_start + header_h + Inches(0.08) + ri * row_h + int(row_h / 2) - Inches(0.08)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx), int(cy), Inches(0.16), Inches(0.16))
            dot.fill.solid()
            dot.fill.fore_color.rgb = hex_to_rgb(tokens['colors']['accent1'])
            dot.line.fill.background()


def render_two_column(slide, element, tokens):
    col_w = int((CONTENT_W - Inches(0.30)) / 2)
    divider_x = int(MARGIN_L) + col_w + int(Inches(0.15))

    # Vertical divider line
    add_vline(slide, divider_x, int(CONTENT_TOP) + Inches(0.10),
             int(CONTENT_H) - Inches(0.20), hex_to_rgb(tokens['colors']['accent1']), Pt(1))

    for idx, side in enumerate(['left', 'right']):
        data = element.get(side, {})
        if not data: continue
        x = int(MARGIN_L) + idx * (col_w + Inches(0.30))

        # Heading
        add_text_box(slide, x, int(CONTENT_TOP) + Inches(0.10), col_w, Inches(0.40),
                    data.get('heading', ''), tokens['fonts']['heading'], Pt(14),
                    hex_to_rgb(tokens['colors']['dk1']), bold=True)

        # Separator under heading
        add_hline(slide, x, int(CONTENT_TOP) + Inches(0.55),
                 col_w, hex_to_rgb(tokens['colors']['accent1']), Pt(1))

        # Body
        add_text_box(slide, x, int(CONTENT_TOP) + Inches(0.70), col_w, int(CONTENT_H) - Inches(0.90),
                    data.get('body', ''), tokens['fonts']['body'], Pt(11),
                    hex_to_rgb(tokens['colors']['dk2']))


def render_stats_row(slide, element, tokens):
    """Massive typography infographic stats blocks, vertically and horizontally centered."""
    items = element.get('items', [])
    if not items: return
    n = len(items)
    gap = Inches(0.30)
    item_w = int((CONTENT_W - gap * (n - 1)) / n)
    card_h = int(CONTENT_H * 0.65)
    mid_y = int(CONTENT_TOP) + int((CONTENT_H - card_h) / 2)

    for i, item in enumerate(items):
        x = int(MARGIN_L) + i * (item_w + int(gap))

        # Add tinted backing card
        accent_hex = tokens['colors']['accent1']
        add_card(slide, x, mid_y, item_w, card_h,
                hex_to_rgb(accent_hex),
                border_color=hex_to_rgb(accent_hex), border_w=Pt(1.5), alpha=10000, shadow=True)

        # Big number - centered
        add_text_box(slide, x, mid_y + Inches(0.2), item_w, Inches(1.2),
                    item.get('value', ''), tokens['fonts']['heading'], Pt(42),
                    hex_to_rgb(accent_hex), bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Accent thin underline under number
        add_hline(slide, x + int(item_w * 0.25), mid_y + Inches(1.50),
                 int(item_w * 0.50), hex_to_rgb(accent_hex), Pt(2))

        # Label - vertically distributed to fill remaining space
        add_text_box(slide, x + Inches(0.15), mid_y + Inches(1.6), item_w - Inches(0.3), card_h - Inches(1.8),
                    item.get('label', ''), tokens['fonts']['body'], Pt(12),
                    hex_to_rgb(tokens['colors']['dk1']), bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def render_quote(slide, element, tokens):
    cw = int(CONTENT_W)
    # Opening quote
    add_text_box(slide, MARGIN_L + Inches(0.50), CONTENT_TOP + Inches(0.20),
                Inches(0.60), Inches(0.60),
                '\u201C', tokens['fonts']['heading'], Pt(48),
                hex_to_rgb(tokens['colors']['dk1']), bold=True)

    # Quote text
    add_text_box(slide, MARGIN_L + Inches(0.80), CONTENT_TOP + Inches(0.60),
                cw - Inches(1.60), Inches(2.50),
                element.get('quote', ''), tokens['fonts']['body'], Pt(14),
                hex_to_rgb(tokens['colors']['dk1']),
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP)

    # Separator
    add_hline(slide, int(MARGIN_L) + Inches(0.80), int(CONTENT_TOP) + Inches(3.30),
             Inches(2.0), hex_to_rgb(tokens['colors']['accent1']), Pt(1.5))

    # Attribution
    add_text_box(slide, MARGIN_L + Inches(0.80), CONTENT_TOP + Inches(3.50),
                cw - Inches(1.60), Inches(0.40),
                f"\u2014 {element.get('attribution', '')}", tokens['fonts']['body'], Pt(11),
                hex_to_rgb(tokens['colors']['dk2']), bold=True)


def render_comparison(slide, element, tokens):
    """High-contrast split-screen comparison block."""
    col_w = int((CONTENT_W - Inches(0.40)) / 2)
    divider_x = int(MARGIN_L) + col_w + int(Inches(0.20))

    for idx, side in enumerate(['left', 'right']):
        data = element.get(side, {})
        if not data: continue
        x = int(MARGIN_L) + idx * (col_w + Inches(0.40))
        # High contrast: Use primary accent for Left, secondary for Right
        accent_hex = tokens['colors']['accent1'] if idx == 0 else tokens['colors'].get('accent3', tokens['colors']['accent2'])

        # Heavy Header Strip
        strip_h = Inches(0.60)
        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, int(CONTENT_TOP), col_w, int(strip_h))
        strip.fill.solid()
        strip.fill.fore_color.rgb = hex_to_rgb(accent_hex)
        strip.line.fill.background()
        
        add_text_box(slide, x + Inches(0.10), int(CONTENT_TOP), col_w - Inches(0.20), strip_h,
                    data.get('title', '').upper(), tokens['fonts']['heading'], Pt(14),
                    hex_to_rgb(tokens['colors']['lt1']), bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Tinted Content Background Drop
        card_h = int(CONTENT_H) - int(strip_h) - Inches(0.10)
        add_card(slide, x, int(CONTENT_TOP) + int(strip_h) + Inches(0.10), col_w, card_h,
                hex_to_rgb(accent_hex), border_color=None, alpha=10000, shadow=True)

        # Bullet List Points
        points = data.get('points', [])
        if points:
            _add_bullet_list(slide, x + Inches(0.20), int(CONTENT_TOP) + int(strip_h) + Inches(0.30),
                           col_w - Inches(0.40), card_h - Inches(0.40),
                           points, tokens, Pt(11))


def add_vector_badge(slide, x, y, size, icon, bg_color, text_color, font_name, alpha_val='100000'):
    # Base circle badge
    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(size), int(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    if alpha_val != '100000':
        a = OxmlElement('a:alpha')
        a.set('val', str(alpha_val)) # Opacity
        bg.fill.fore_color._color._xClr.append(a)
    
    # Text icon inside, perfectly centered
    tf = bg.text_frame
    _zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = icon
    p.font.name = font_name
    p.font.size = Pt(int((size / 914400) * 32)) # Massive icon relative to badge size
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return bg

def render_icon_grid(slide, element, tokens):
    """Icon grid mapping heavily centered icons and typography with distinct contrast borders."""
    items = element.get('items', [])
    cols = min(element.get('columns', 3), 4)
    if not items: return

    gap = Inches(0.25)
    rows = -(-len(items) // cols)
    row_gap = Inches(0.35)
    card_w = int((CONTENT_W - gap * (cols - 1)) / cols)
    card_h = int((CONTENT_H - row_gap * (rows - 1)) / rows)
    card_h = min(card_h, int(CONTENT_H * 0.8))

    badge_size = Inches(0.60)

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        x = int(MARGIN_L) + col * (card_w + int(gap))
        y = int(CONTENT_TOP) + row * (card_h + int(row_gap))

        # Main tinted transparent card
        accent = tokens['colors']['accent1']
        add_card(slide, x, y, card_w, card_h,
                hex_to_rgb(accent),
                hex_to_rgb(accent), Pt(1.5), alpha=15000, shadow=True)

        # Center-top overlapping Vector Icon Badge (prominent, 100% opacity solid fill)
        badge_y = y - int(badge_size / 2)
        badge_x = x + int(card_w / 2) - int(badge_size / 2)
        add_vector_badge(slide, badge_x, badge_y, badge_size,
                        item.get('icon', '🔹'), hex_to_rgb(tokens['colors']['accent1']),
                        hex_to_rgb(tokens['colors']['lt1']), tokens['fonts']['body'])

        # Bold Centered Title
        title_y = y + int(badge_size/2) + Inches(0.10)
        title_h = Inches(0.40)
        add_text_box(slide, x + Inches(0.10), title_y,
                    card_w - Inches(0.20), title_h,
                    item.get('title', ''), tokens['fonts']['heading'], Pt(12),
                    hex_to_rgb(tokens['colors']['dk1']), bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Accent H-Line
        add_hline(slide, x + int((card_w - Inches(1.0))/2), title_y + title_h + Inches(0.05),
                 Inches(1.0), hex_to_rgb(tokens['colors'].get('accent2', accent)), Pt(1.5))

        # Description (Middle aligned inside remaining space)
        desc_y = title_y + title_h + Inches(0.15)
        add_text_box(slide, x + Inches(0.15), desc_y,
                    card_w - Inches(0.30), card_h - (desc_y - y) - Inches(0.10),
                    item.get('description', ''), tokens['fonts']['body'], Pt(10),
                    hex_to_rgb(tokens['colors']['dk2']),
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)


def render_funnel(slide, element, tokens):
    steps = element.get('steps', [])
    if not steps: return
    n = len(steps)
    step_h = int(CONTENT_H / n)
    max_w = int(CONTENT_W) - Inches(1.0)
    min_w = int(max_w * 0.3)
    colors = [tokens['colors']['accent1'], tokens['colors']['accent2'], tokens['colors']['accent3'],
              tokens['colors']['accent4'], tokens['colors']['accent5']]

    for i, step in enumerate(steps):
        ratio = 1 - (i / max(n - 1, 1)) * 0.7
        w = int(min_w + (max_w - min_w) * ratio)
        x = int(MARGIN_L) + int((int(CONTENT_W) - w) / 2)
        y = int(CONTENT_TOP) + i * step_h

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, step_h - Inches(0.06))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(colors[i % len(colors)])
        shape.line.fill.background()

        label = f"{step.get('label', '')} ({step.get('value', '')})" if step.get('value') else step.get('label', '')
        tf = shape.text_frame
        _zero_margins(tf)
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.name = tokens['fonts']['heading']
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = hex_to_rgb(tokens['colors']['lt1'])
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def render_pyramid(slide, element, tokens):
    levels = element.get('levels', [])
    if not levels: return
    n = len(levels)
    step_h = int(CONTENT_H / n)
    max_w = int(CONTENT_W) - Inches(1.0)
    min_w = int(max_w * 0.25)
    colors = [tokens['colors']['accent1'], tokens['colors']['accent2'],
              tokens['colors']['accent3'], tokens['colors']['accent4']]

    for i, level in enumerate(levels):
        ratio = i / max(n - 1, 1)
        w = int(min_w + (max_w - min_w) * ratio)
        x = int(MARGIN_L) + int((int(CONTENT_W) - w) / 2)
        y = int(CONTENT_TOP) + i * step_h

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, step_h - Inches(0.06))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(colors[i % len(colors)])
        shape.line.fill.background()

        txt = f"{level.get('label', '')}: {level.get('description', '')}" if level.get('description') else level.get('label', '')
        tf = shape.text_frame
        _zero_margins(tf)
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].text = txt
        tf.paragraphs[0].font.name = tokens['fonts']['heading']
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = hex_to_rgb(tokens['colors']['lt1'])
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def render_matrix(slide, element, tokens):
    """High-contrast 2x2 matrix map using heavy boundary lines and completely filled alpha masks."""
    quadrants = element.get('quadrants', [])
    if len(quadrants) < 4: return

    mat_w = int((int(CONTENT_W) - Inches(0.15)) / 2)
    mat_h = int((int(CONTENT_H) - Inches(0.15)) / 2)
    colors = [tokens['colors']['accent1'], tokens['colors'].get('accent3', tokens['colors']['accent2']),
              tokens['colors']['accent2'], tokens['colors']['accent4']]
    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]

    for idx, (q, (c, r)) in enumerate(zip(quadrants, positions)):
        x = int(MARGIN_L) + c * (mat_w + Inches(0.15))
        y = int(CONTENT_TOP) + r * (mat_h + Inches(0.15))

        # Matrix block background (15% filled assigned color)
        add_card(slide, x, y, mat_w, mat_h,
                hex_to_rgb(colors[idx]),
                border_color=hex_to_rgb(colors[idx]), border_w=Pt(2), alpha=15000, shadow=True)

        # Full width header within block
        header_h = Inches(0.40)
        h_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, mat_w, int(header_h))
        h_strip.fill.solid()
        h_strip.fill.fore_color.rgb = hex_to_rgb(colors[idx])
        h_strip.line.fill.background()

        add_text_box(slide, x + Inches(0.10), y, mat_w - Inches(0.20), header_h,
                    q.get('label', '').upper(), tokens['fonts']['heading'], Pt(11),
                    hex_to_rgb(tokens['colors']['lt1']), bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        items = q.get('items', [])
        if items:
            _add_bullet_list(slide, x + Inches(0.15), y + header_h + Inches(0.10),
                           mat_w - Inches(0.30), mat_h - header_h - Inches(0.20),
                           items, tokens, Pt(10))


def render_swot(slide, element, tokens):
    render_matrix(slide, {
        'quadrants': [
            {'label': 'Strengths', 'items': element.get('strengths', [])},
            {'label': 'Weaknesses', 'items': element.get('weaknesses', [])},
            {'label': 'Opportunities', 'items': element.get('opportunities', [])},
            {'label': 'Threats', 'items': element.get('threats', [])},
        ]
    }, tokens)


def render_cycle(slide, element, tokens):
    steps = element.get('steps', [])
    if not steps: return
    n = len(steps)
    center_x = int(MARGIN_L) + int(CONTENT_W / 2)
    center_y = int(CONTENT_TOP) + int(CONTENT_H / 2)
    radius = min(int(CONTENT_W), int(CONTENT_H)) / 2 - Inches(0.70)
    node_w = Inches(1.80)
    node_h = Inches(0.80)
    colors = [tokens['colors']['accent1'], tokens['colors']['accent2'], tokens['colors']['accent3'],
              tokens['colors']['accent4'], tokens['colors']['accent5']]

    for i, step in enumerate(steps):
        angle = (2 * math.pi * i) / n - math.pi / 2
        cx = int(center_x + radius * math.cos(angle) - node_w / 2)
        cy = int(center_y + radius * math.sin(angle) - node_h / 2)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, int(node_w), int(node_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(colors[i % len(colors)])
        shape.line.fill.background()

        tf = shape.text_frame
        _zero_margins(tf)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = step.get('title', '')
        run.font.name = tokens['fonts']['heading']
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(tokens['colors']['lt1'])
        p.alignment = PP_ALIGN.CENTER

        if step.get('description'):
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = step['description']
            r2.font.name = tokens['fonts']['body']
            r2.font.size = Pt(8)
            r2.font.color.rgb = hex_to_rgb(tokens['colors']['lt1'])
            p2.alignment = PP_ALIGN.CENTER


def render_gauge(slide, element, tokens):
    center_x = int(MARGIN_L) + int(CONTENT_W / 2)
    center_y = int(CONTENT_TOP) + int(CONTENT_H * 0.45)

    value_str = f"{element.get('value', 0)}{element.get('unit', '%')}"
    add_text_box(slide, center_x - Inches(1.5), center_y - Inches(0.35), Inches(3.0), Inches(0.70),
                value_str, tokens['fonts']['heading'], Pt(40),
                hex_to_rgb(tokens['colors']['dk1']), bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, center_x - Inches(2.0), center_y + Inches(0.50), Inches(4.0), Inches(0.40),
                element.get('label', ''), tokens['fonts']['body'], Pt(14),
                hex_to_rgb(tokens['colors']['dk2']),
                align=PP_ALIGN.CENTER)

    add_hline(slide, center_x - Inches(1.0), center_y + Inches(1.0),
             Inches(2.0), hex_to_rgb(tokens['colors']['accent1']), Pt(2))


def render_kpi_cards(slide, element, tokens):
    items = element.get('items', [])
    if not items: return
    n = len(items)
    gap = Inches(0.20)
    card_w = int((CONTENT_W - gap * (n - 1)) / n)
    card_h = int(CONTENT_H * 0.70)
    mid_y = int(CONTENT_TOP) + int((CONTENT_H - card_h) / 2)

    for i, item in enumerate(items):
        x = int(MARGIN_L) + i * (card_w + int(gap))
        add_card(slide, x, mid_y, card_w, int(card_h),
                hex_to_rgb(tokens['colors']['lt2']),
                hex_to_rgb(tokens['colors']['accent1']), Pt(0))

        # Value
        add_text_box(slide, x, mid_y + Inches(0.4), card_w, Inches(1.0),
                    item.get('value', ''), tokens['fonts']['heading'], Pt(36),
                    hex_to_rgb(tokens['colors']['dk1']), bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Separator
        add_hline(slide, x + Inches(0.3), mid_y + Inches(1.6),
                 card_w - Inches(0.6), hex_to_rgb(tokens['colors']['accent1']), Pt(1.5))

        # Label
        add_text_box(slide, x + Inches(0.2), mid_y + Inches(1.8), card_w - Inches(0.4), Inches(0.6),
                    item.get('label', ''), tokens['fonts']['body'], Pt(12),
                    hex_to_rgb(tokens['colors']['dk2']),
                    align=PP_ALIGN.CENTER)

        # Trend
        if item.get('change'):
            arrow = '\u25B2' if item.get('trend') == 'up' else '\u25BC' if item.get('trend') == 'down' else '\u25C6'
            trend_color = tokens['colors'].get('accent3', tokens['colors'].get('accent6', tokens['colors']['accent2'])) if item.get('trend') == 'up' else tokens['colors']['accent1']
            add_text_box(slide, x, mid_y + Inches(2.6), card_w, Inches(0.5),
                        f"{arrow} {item['change']}", tokens['fonts']['body'], Pt(14),
                        hex_to_rgb(trend_color), bold=True,
                        align=PP_ALIGN.CENTER)


import os
def add_image_placeholder(slide, x, y, w, h, query, tokens, image_url=None, show_icon_text=True):
    if image_url:
        img_path = image_url.strip()
        base_dirs = [".", os.path.dirname(r"Sample Files\Sample Files\Accenture Tech Acquisition Analysis\Accenture Tech Acquisition Analysis.md")]
        for base in base_dirs:
            full_test = os.path.normpath(os.path.join(base, img_path))
            if os.path.exists(full_test):
                try:
                    slide.shapes.add_picture(full_test, int(x), int(y), int(w), int(h))
                    return
                except Exception:
                    pass

    bg = add_card(slide, x, y, w, h, hex_to_rgb(tokens['colors']['lt2']), hex_to_rgb(tokens['colors']['accent2']), Pt(1), alpha=50000)
    
    if show_icon_text:
        add_text_box(slide, x, y + int(h/2) - Inches(0.4), w, Inches(0.5),
                    "🖼️", tokens['fonts']['body'], Pt(32), hex_to_rgb(tokens['colors']['dk2']),
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        add_text_box(slide, x, y + int(h/2) + Inches(0.3), w, Inches(0.3),
                    f"Image: {query}" if query else "Image Placeholder", tokens['fonts']['body'], Pt(11), 
                    hex_to_rgb(tokens['colors']['dk2']), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

def render_image_text(slide, element, tokens):
    col_w = int((CONTENT_W - Inches(0.30)) / 2)
    content = element.get('content', {})
    img_left = content.get('image_side', 'right') == 'left'
    query = content.get('image_query')

    image_url = content.get('image_url')

    text_x = int(MARGIN_L) + (col_w + Inches(0.30) if img_left else 0)
    img_x = int(MARGIN_L) + (0 if img_left else col_w + Inches(0.30))

    # Image placeholder or Actual Image
    add_image_placeholder(slide, img_x, int(CONTENT_TOP), col_w, int(CONTENT_H), query, tokens, image_url)

    # Text side
    add_text_box(slide, text_x, int(CONTENT_TOP) + Inches(0.15), col_w, Inches(0.40),
                content.get('heading', ''), tokens['fonts']['heading'], Pt(16),
                hex_to_rgb(tokens['colors']['dk1']), bold=True)

    add_hline(slide, text_x, int(CONTENT_TOP) + Inches(0.60),
             col_w, hex_to_rgb(tokens['colors']['accent1']), Pt(1))

    add_text_box(slide, text_x, int(CONTENT_TOP) + Inches(0.75), col_w, int(CONTENT_H) - Inches(1.0),
                content.get('body', ''), tokens['fonts']['body'], Pt(11),
                hex_to_rgb(tokens['colors']['dk2']))


def render_waterfall(slide, element, tokens):
    steps = element.get('steps', [])
    if not steps: return
    n = len(steps)
    gap = Inches(0.08)
    bar_w = int((CONTENT_W - gap * (n - 1)) / n)
    max_val = max(abs(s.get('value', 0)) for s in steps) or 1
    base_y = int(CONTENT_TOP) + int(CONTENT_H) - Inches(0.70)

    cumulative = 0
    for i, step in enumerate(steps):
        x = int(MARGIN_L) + i * (bar_w + int(gap))
        val = step.get('value', 0)
        bar_h = max(int(abs(val) / max_val * int(CONTENT_H) * 0.4), Inches(0.15))
        is_total = step.get('is_total', False)

        if is_total:
            y = base_y - bar_h
        elif val >= 0:
            y = base_y - int((cumulative + val) / max_val * int(CONTENT_H) * 0.4)
        else:
            y = base_y - int(cumulative / max_val * int(CONTENT_H) * 0.4)

        fill_color = tokens['colors']['accent1'] if is_total else (
            tokens['colors'].get('accent3', tokens['colors'].get('accent6', tokens['colors']['accent2'])) if val >= 0 else tokens['colors']['accent1'])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, bar_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        shape.line.fill.background()

        add_text_box(slide, x, base_y + Inches(0.05), bar_w, Inches(0.30),
                    step.get('label', ''), tokens['fonts']['body'], Pt(8),
                    hex_to_rgb(tokens['colors']['dk2']),
                    align=PP_ALIGN.CENTER)

        if not is_total: cumulative += val


# ─── Native Placeholder Rendering (Hybrid) ────────────────────────────────────

def populate_placeholders(slide, slide_def, tokens):
    """
    Attempt to render content natively using template placeholders.
    Returns True if successful, False if we need to fall back to the shape engine.
    """
    # 1. Discover placeholders
    placeholders = {'title': [], 'body': [], 'picture': []}
    for ph in slide.placeholders:
        ph_type = str(ph.placeholder_format.type).split('(')[0].strip().lower().replace(' ', '_')
        if ph_type in ('title', 'center_title', 'vertical_title'):
            placeholders['title'].append(ph)
        elif ph_type in ('body', 'object', 'vertical_body'):
            placeholders['body'].append(ph)
        elif ph_type in ('picture', 'bitmap'):
            placeholders['picture'].append(ph)

    for k in placeholders:
        placeholders[k].sort(key=lambda p: (p.top, p.left))

    elements = slide_def.get('elements', [])
    complex_types = {
        'hero', 'icon_grid', 'pyramid', 'funnel', 'matrix', 'swot', 
        'waterfall', 'cycle', 'gauge', 'kpi_cards', 'chart', 'table',
        'timeline'
    }

    if elements:
        first_el = elements[0]
        # Skip if complex visual
        if first_el.get('type') in complex_types:
            return False
            
        # Two column needs at least 2 body placeholders
        if first_el.get('type') == 'two_column' and len(placeholders['body']) < 2:
            return False

        # If it's a simple type but we have 0 body placeholders, we must fallback
        if not placeholders['body']:
            return False

    # 2. Populate Title
    title_text = slide_def.get('title', '')
    if title_text and placeholders['title']:
        try:
            from content_fitter import calculate_fit_font_size
            title_ph = placeholders['title'][0]
            title_ph.text = title_text
            fit_sz = calculate_fit_font_size(title_text, title_ph.width/914400, title_ph.height/914400, max_font_pt=32) if title_ph.width else 24
            
            for p in title_ph.text_frame.paragraphs:
                p.font.name = tokens['fonts']['heading']
                p.font.bold = True
                p.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])
                p.font.size = Pt(fit_sz)
        except Exception as e:
            print(f"Error fitting title: {e}")

    # 3. Populate Body
    if elements and placeholders['body']:
        first_el = elements[0]
        etype = first_el.get('type', '')
        
        try:
            from content_fitter import calculate_bullet_fit
            
            if etype == 'bullets':
                items = first_el.get('items', [])
                body_ph = placeholders['body'][0]
                fit_sz = calculate_bullet_fit(items, body_ph.width/914400, body_ph.height/914400, max_font_pt=18) if body_ph.width else 14
                    
                tf = body_ph.text_frame
                tf.clear()
                
                for i, item in enumerate(items):
                    text = item if isinstance(item, str) else item.get('text', '')
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = text
                    p.font.name = tokens['fonts']['body']
                    p.font.size = Pt(fit_sz)
                    p.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])
                    p.level = 0
            
            elif etype == 'two_column':
                cols = first_el.get('columns', [])
                for col_idx, col in enumerate(cols[:2]):
                    items = col.get('items', [])
                    body_ph = placeholders['body'][col_idx]
                    fit_sz = calculate_bullet_fit(items, body_ph.width/914400, body_ph.height/914400, max_font_pt=16) if body_ph.width else 14
                        
                    tf = body_ph.text_frame
                    tf.clear()
                    
                    header = col.get('header', '')
                    start_idx = 0
                    if header:
                        p = tf.paragraphs[0]
                        p.text = header
                        p.font.name = tokens['fonts']['heading']
                        p.font.size = Pt(fit_sz + 2)
                        p.font.bold = True
                        p.font.color.rgb = hex_to_rgb(tokens['colors']['accent1'])
                        start_idx = 1
                    
                    for i, item in enumerate(items):
                        text = item if isinstance(item, str) else item.get('text', '')
                        p = tf.paragraphs[0] if (i == 0 and start_idx == 0) else tf.add_paragraph()
                        p.text = text
                        p.font.name = tokens['fonts']['body']
                        p.font.size = Pt(fit_sz)
                        p.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])
                        p.level = 0
            else:
                return False
        except Exception as e:
            print(f"Error in native generation: {e}")
            return False

    return True


# ─── Renderer Dispatch ────────────────────────────────────────────────────────

RENDERERS = {
    'grid': render_grid,
    'timeline': render_timeline,
    'hero': render_hero,
    'bullets': render_bullets,
    'chart': render_chart,
    'table': render_table,
    'two_column': render_two_column,
    'stats_row': render_stats_row,
    'quote': render_quote,
    'comparison': render_comparison,
    'icon_grid': render_icon_grid,
    'waterfall': render_waterfall,
    'funnel': render_funnel,
    'pyramid': render_pyramid,
    'matrix': render_matrix,
    'swot': render_swot,
    'cycle': render_cycle,
    'gauge': render_gauge,
    'kpi_cards': render_kpi_cards,
    'image_text': render_image_text,
}


# ─── Main Compile ─────────────────────────────────────────────────────────────

def compile_presentation(tokens_path, plan_path, template_path, output_path):
    with open(tokens_path, 'r', encoding='utf-8') as f:
        tokens = json.load(f)
    with open(plan_path, 'r', encoding='utf-8') as f:
        plan = json.load(f)

    prs = Presentation(template_path)

    # Strip "Source" text boxes from the master layouts directly
    for layout in prs.slide_layouts:
        for shape in list(layout.shapes):
            if shape.has_text_frame:
                try:
                    if shape.text.strip().startswith('Source'):
                        element = shape.element
                        element.getparent().remove(element)
                except Exception:
                    pass

    # Remove existing template slides
    xml_slides = prs.slides._sldIdLst
    for s in list(xml_slides):
        prs.part.drop_rel(s.rId)
        xml_slides.remove(s)

    report_name = tokens.get('template_name', '')
    slides = plan.get('slides', [])
    section_counter = 0

    print(f"Compiling {len(slides)} slides from template: {template_path}")

    for idx, slide_def in enumerate(slides):
        layout_type = slide_def.get('layout', 'content')
        layout = get_layout(prs, layout_type)
        slide = prs.slides.add_slide(layout)

        print(f"  Slide {idx+1}: [{layout_type}] {slide_def.get('title', '')} (layout: {layout.name})")

        # Remove "Source" text box from template if present
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                try:
                    if shape.text.strip().startswith('Source'):
                        element = shape.element
                        element.getparent().remove(element)
                except Exception:
                    pass

        # ─── Cover ─────────────────────────────────────────────
        if layout_type == 'cover':
            placeholders = sorted(list(slide.placeholders), key=lambda x: x.top)
            if len(placeholders) >= 2:
                placeholders[0].text = slide_def.get('title', '')
                for p in placeholders[0].text_frame.paragraphs:
                    p.font.name = tokens['fonts']['heading']
                    p.font.size = Pt(36)
                    p.font.bold = True
                    p.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])

                placeholders[1].text = slide_def.get('subtitle', '')
                for p in placeholders[1].text_frame.paragraphs:
                    p.font.name = tokens['fonts']['body']
                    p.font.size = Pt(14)
                    p.font.color.rgb = hex_to_rgb(tokens['colors']['dk2'])
            elif len(placeholders) == 1:
                placeholders[0].text = slide_def.get('title', '')
            continue

        # ─── Thank You ─────────────────────────────────────────
        if layout_type == 'thank_you':
            if slide.shapes.title is not None:
                slide.shapes.title.text = slide_def.get('title', 'Thank You')
                slide.shapes.title.text_frame.word_wrap = True
            add_footer(slide, tokens, idx + 1, report_name)
            continue

        # ─── Divider ───────────────────────────────────────────
        if layout_type == 'divider':
            placeholders = list(slide.placeholders)
            if placeholders:
                placeholders[0].text = slide_def.get('title', '')
                for p in placeholders[0].text_frame.paragraphs:
                    p.font.name = tokens['fonts']['heading']
                    p.font.size = Pt(28)
                    p.font.bold = True
                    p.font.color.rgb = hex_to_rgb(tokens['colors']['dk1'])
            add_footer(slide, tokens, idx + 1, report_name)
            continue

        # ─── Content / Chart ───────────────────────────────────
        section_counter += 1

        # Section label
        add_section_label(slide, f"SECTION {section_counter}", tokens)

        # Footer
        add_footer(slide, tokens, idx + 1, report_name)

        # Try native placeholder filling first (Hybrid Approach)
        is_native = False
        elements = slide_def.get('elements', [])
        
        if elements:
            etype = elements[0].get('type', '')
            if etype in ('bullets', 'two_column'):
                is_native = populate_placeholders(slide, slide_def, tokens)
                if is_native:
                    print(f"    -> Rendered natively using template placeholders")

        if not is_native:
            # Fall back to custom shape rendering engine
            add_title(slide, slide_def.get('title', ''), tokens)
            
            for element in slide_def.get('elements', []):
                etype = element.get('type', '')
                renderer = RENDERERS.get(etype)
                if renderer:
                    try:
                        renderer(slide, element, tokens)
                    except Exception as e:
                        print(f"    WARNING: Error rendering '{etype}': {e}")
                else:
                    print(f"    WARNING: Unknown element type '{etype}', skipping.")

    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")


# ─── CLI ──────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Stage 3 -- Compile: Tokens + Plan -> .pptx')
    parser.add_argument('tokens', help='Path to design_tokens.json')
    parser.add_argument('plan', help='Path to ui_plan.json')
    parser.add_argument('template', help='Path to the original .pptx template')
    parser.add_argument('output', help='Path to output .pptx file')
    args = parser.parse_args()

    compile_presentation(args.tokens, args.plan, args.template, args.output)
