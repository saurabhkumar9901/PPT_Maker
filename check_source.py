import pptx
prs = pptx.Presentation(r'Slide Master\Slide Master\Template_Accenture Tech Acquisition Analysis.pptx')
found = False
for master in prs.slide_masters:
    for shape in master.shapes:
        if shape.has_text_frame and 'Source' in shape.text:
            print(f'Found in master: {shape.text}')
            found = True
for layout in prs.slide_layouts:
    for shape in layout.shapes:
        if shape.has_text_frame and 'Source' in shape.text:
            print(f'Found in layout {layout.name}: {shape.text[:20]}')
            found = True
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and 'Source' in shape.text:
            print(f'Found in slide: {shape.text}')
            found = True
if not found:
    print('Not found')
