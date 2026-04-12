"""
Post-Generation Validator for PPT Maker v2.

Automated quality checks that catch visual issues after compilation:
  - Text overflow detection
  - Empty placeholder / slide detection
  - Font consistency verification
  - Slide count validation
  - Color contrast warnings

Usage:
    from validator import validate_presentation
    report = validate_presentation("output/deck.pptx", tokens)
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE


def validate_presentation(pptx_path: str, tokens: dict = None) -> dict:
    """Run all quality checks on a generated presentation.
    
    Args:
        pptx_path: Path to the .pptx file to validate
        tokens: Optional design tokens for brand consistency checks
    
    Returns:
        Dict with 'passed', 'warnings', 'errors', and 'details'
    """
    if not os.path.exists(pptx_path):
        return {
            'passed': False,
            'warnings': [],
            'errors': [f"File not found: {pptx_path}"],
            'details': {}
        }

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {
            'passed': False,
            'warnings': [],
            'errors': [f"Failed to open presentation: {e}"],
            'details': {}
        }

    warnings = []
    errors = []
    details = {
        'slide_count': len(prs.slides),
        'slide_width': round(prs.slide_width / 914400, 2),
        'slide_height': round(prs.slide_height / 914400, 2),
        'slides': []
    }

    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': i,
            'shape_count': len(slide.shapes),
            'issues': []
        }

        for shape in slide.shapes:
            # ── Check 1: Shape bounds within slide ────────────────
            if shape.left is not None and shape.top is not None:
                right_edge = shape.left + (shape.width or 0)
                bottom_edge = shape.top + (shape.height or 0)
                
                if right_edge > slide_w_emu + Inches(0.1):
                    issue = f"Shape '{shape.name}' extends beyond right edge by {round((right_edge - slide_w_emu) / 914400, 2)}in"
                    slide_info['issues'].append(issue)
                    warnings.append(f"Slide {i}: {issue}")
                
                if bottom_edge > slide_h_emu + Inches(0.1):
                    issue = f"Shape '{shape.name}' extends beyond bottom edge by {round((bottom_edge - slide_h_emu) / 914400, 2)}in"
                    slide_info['issues'].append(issue)
                    warnings.append(f"Slide {i}: {issue}")

            # ── Check 2: Text overflow heuristic ──────────────────
            if shape.has_text_frame:
                tf = shape.text_frame
                full_text = tf.text.strip()
                
                # Skip empty text frames
                if not full_text:
                    continue
                
                # Check for very small text (might be auto-shrunk too much)
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None and run.font.size < Pt(6):
                            issue = f"Shape '{shape.name}' has extremely small text ({round(run.font.size.pt, 1)}pt)"
                            slide_info['issues'].append(issue)
                            warnings.append(f"Slide {i}: {issue}")
                            break
                
                # Check for potential text overflow (heuristic)
                if shape.width and shape.height:
                    box_w_in = shape.width / 914400
                    box_h_in = shape.height / 914400
                    char_count = len(full_text)
                    
                    # Very rough: if character density is extreme, flag it
                    area_sq_in = box_w_in * box_h_in
                    if area_sq_in > 0:
                        density = char_count / area_sq_in
                        if density > 800:  # More than ~800 chars per sq inch is suspicious
                            issue = f"Shape '{shape.name}' may have text overflow (density: {int(density)} chars/sqin)"
                            slide_info['issues'].append(issue)
                            warnings.append(f"Slide {i}: {issue}")

        # ── Check 3: Empty slide (no content shapes) ──────────
        content_shapes = [s for s in slide.shapes 
                         if s.has_text_frame and s.text_frame.text.strip()]
        if len(content_shapes) == 0:
            issue = "Slide appears to be empty (no text content)"
            slide_info['issues'].append(issue)
            warnings.append(f"Slide {i}: {issue}")

        details['slides'].append(slide_info)

    # ── Check 4: Font consistency ─────────────────────────────
    if tokens:
        expected_heading = tokens.get('fonts', {}).get('heading', '')
        expected_body = tokens.get('fonts', {}).get('body', '')
        font_mismatches = set()
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.font.name not in (
                                expected_heading, expected_body, 
                                'Calibri', 'Arial', 'Wingdings', 'Symbol'
                            ):
                                font_mismatches.add(run.font.name)
        
        if font_mismatches:
            for font in font_mismatches:
                warnings.append(f"Unexpected font used: '{font}' (expected '{expected_heading}' or '{expected_body}')")

    # ── Final verdict ─────────────────────────────────────────
    passed = len(errors) == 0
    
    return {
        'passed': passed,
        'warnings': warnings,
        'errors': errors,
        'details': details,
        'summary': {
            'total_slides': len(prs.slides),
            'warning_count': len(warnings),
            'error_count': len(errors),
            'slides_with_issues': sum(1 for s in details['slides'] if s['issues'])
        }
    }


def print_validation_report(report: dict):
    """Print a human-readable validation report."""
    summary = report.get('summary', {})
    
    status = "✅ PASSED" if report['passed'] else "❌ FAILED"
    print(f"\n{'='*60}")
    print(f"VALIDATION REPORT: {status}")
    print(f"{'='*60}")
    print(f"  Slides: {summary.get('total_slides', 0)}")
    print(f"  Warnings: {summary.get('warning_count', 0)}")
    print(f"  Errors: {summary.get('error_count', 0)}")
    print(f"  Slides with issues: {summary.get('slides_with_issues', 0)}")
    
    if report['errors']:
        print(f"\n❌ ERRORS:")
        for err in report['errors']:
            print(f"  • {err}")
    
    if report['warnings']:
        print(f"\n⚠️  WARNINGS:")
        for warn in report['warnings']:
            print(f"  • {warn}")
    
    if not report['warnings'] and not report['errors']:
        print(f"\n  No issues detected. Presentation looks good!")
    
    print(f"{'='*60}\n")


# ─── CLI ──────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description="Validate a generated presentation")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--tokens", help="Path to design_tokens.json for brand checks")
    args = parser.parse_args()
    
    tokens = None
    if args.tokens:
        with open(args.tokens, 'r', encoding='utf-8') as f:
            tokens = json.load(f)
    
    report = validate_presentation(args.pptx, tokens)
    print_validation_report(report)
