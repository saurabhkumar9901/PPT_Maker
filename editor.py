"""
Editor for PPT Maker v2 — Surgical edits to existing presentations.

Provides three modes:
  1. --inventory: List all text shapes in a presentation
  2. --replace: Replace text across slides  
  3. --reorder: Reorder slides by index list

Usage:
    python editor.py deck.pptx --inventory
    python editor.py deck.pptx --replace '{"slide":3,"old":"2025","new":"2026"}'
    python editor.py deck.pptx --reorder "0,2,1,3,4" -o reordered.pptx

Ported and adapted from pptx-from-layouts-skill.
"""

import os
import json
import copy
import argparse
from pptx import Presentation
from pptx.util import Emu
from lxml import etree


def get_inventory(pptx_path: str) -> list:
    """Extract a complete text inventory from a presentation.
    
    Returns a list of dicts, one per text element:
        {slide: int, shape: str, placeholder_idx: int|None, text: str, 
         left: float, top: float, width: float, height: float}
    """
    prs = Presentation(pptx_path)
    inventory = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                entry = {
                    'slide': slide_idx,
                    'shape_name': shape.name,
                    'shape_id': shape.shape_id,
                    'placeholder_idx': None,
                    'text': shape.text_frame.text[:200],  # Truncate for readability
                    'position': {
                        'left': round(shape.left / 914400, 2) if shape.left else None,
                        'top': round(shape.top / 914400, 2) if shape.top else None,
                        'width': round(shape.width / 914400, 2) if shape.width else None,
                        'height': round(shape.height / 914400, 2) if shape.height else None,
                    }
                }
                
                if shape.is_placeholder:
                    entry['placeholder_idx'] = shape.placeholder_format.idx
                
                inventory.append(entry)
    
    return inventory


def print_inventory(inventory: list):
    """Pretty-print a text inventory."""
    print(f"\n{'='*70}")
    print(f"TEXT INVENTORY ({len(inventory)} shapes with text)")
    print(f"{'='*70}")
    
    current_slide = 0
    for item in inventory:
        if item['slide'] != current_slide:
            current_slide = item['slide']
            print(f"\n── Slide {current_slide} ──")
        
        ph_info = f" [PH:{item['placeholder_idx']}]" if item['placeholder_idx'] is not None else ""
        text_preview = item['text'].replace('\n', ' ↵ ')
        if len(text_preview) > 80:
            text_preview = text_preview[:77] + "..."
        
        print(f"  {item['shape_name']}{ph_info}: \"{text_preview}\"")
    
    print(f"\n{'='*70}\n")


def replace_text(pptx_path: str, replacements: list, output_path: str = None) -> dict:
    """Replace text in a presentation.
    
    Args:
        pptx_path: Path to the source .pptx
        replacements: List of replacement dicts, each with:
            - slide (int, optional): Target slide number (1-indexed). If omitted, applies to all.
            - old (str): Text to find
            - new (str): Replacement text
        output_path: Where to save. If None, overwrites the source.
    
    Returns:
        Dict with 'total_replacements' count and 'details'
    """
    prs = Presentation(pptx_path)
    save_path = output_path or pptx_path
    
    total = 0
    details = []
    
    for repl in replacements:
        old_text = repl.get('old', '')
        new_text = repl.get('new', '')
        target_slide = repl.get('slide')
        
        if not old_text:
            continue
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            if target_slide and slide_idx != target_slide:
                continue
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            total += 1
                            details.append({
                                'slide': slide_idx,
                                'shape': shape.name,
                                'old': old_text,
                                'new': new_text,
                            })
    
    os.makedirs(os.path.dirname(save_path) or '.', exist_ok=True)
    prs.save(save_path)
    
    return {
        'total_replacements': total,
        'output_path': save_path,
        'details': details,
    }


def reorder_slides(pptx_path: str, order: list, output_path: str = None) -> str:
    """Reorder slides by index list.
    
    Args:
        pptx_path: Path to the source .pptx
        order: List of 0-indexed slide positions, e.g. [0, 2, 1, 3, 4]
        output_path: Where to save
    
    Returns:
        Path to the saved file
    """
    prs = Presentation(pptx_path)
    save_path = output_path or pptx_path
    
    slides_list = prs.slides._sldIdLst
    slide_elements = list(slides_list)
    
    # Validate indices
    n = len(slide_elements)
    if len(order) != n:
        raise ValueError(f"Order list has {len(order)} items but presentation has {n} slides")
    if set(order) != set(range(n)):
        raise ValueError(f"Order list must contain each index 0..{n-1} exactly once")
    
    # Remove all and re-add in new order
    for el in slide_elements:
        slides_list.remove(el)
    
    for idx in order:
        slides_list.append(slide_elements[idx])
    
    os.makedirs(os.path.dirname(save_path) or '.', exist_ok=True)
    prs.save(save_path)
    return save_path


# ─── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="PPT Maker Editor — Surgical edits to existing presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  List all text shapes:
    python editor.py deck.pptx --inventory

  Replace text:
    python editor.py deck.pptx --replace '{"slide":3,"old":"2025","new":"2026"}'

  Multiple replacements from file:
    python editor.py deck.pptx --replace-file changes.json -o edited.pptx

  Reorder slides:
    python editor.py deck.pptx --reorder "0,2,1,3,4" -o reordered.pptx
        """
    )
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--inventory", action="store_true", help="List all text shapes")
    parser.add_argument("--replace", help="JSON string with replacement: {slide, old, new}")
    parser.add_argument("--replace-file", help="Path to JSON file with array of replacements")
    parser.add_argument("--reorder", help="Comma-separated 0-indexed slide order")
    parser.add_argument("-o", "--output", help="Output path (default: overwrite input)")

    args = parser.parse_args()

    if args.inventory:
        inv = get_inventory(args.pptx)
        print_inventory(inv)
    
    elif args.replace:
        repl = json.loads(args.replace)
        if isinstance(repl, dict):
            repl = [repl]
        result = replace_text(args.pptx, repl, args.output)
        print(f"Replaced {result['total_replacements']} occurrences")
        if result['details']:
            for d in result['details']:
                print(f"  Slide {d['slide']}, {d['shape']}: \"{d['old']}\" → \"{d['new']}\"")
        print(f"Saved to: {result['output_path']}")
    
    elif args.replace_file:
        with open(args.replace_file, 'r', encoding='utf-8') as f:
            replacements = json.load(f)
        result = replace_text(args.pptx, replacements, args.output)
        print(f"Replaced {result['total_replacements']} occurrences")
        print(f"Saved to: {result['output_path']}")
    
    elif args.reorder:
        order = [int(x.strip()) for x in args.reorder.split(',')]
        out = reorder_slides(args.pptx, order, args.output)
        print(f"Slides reordered. Saved to: {out}")
    
    else:
        parser.print_help()


if __name__ == '__main__':
    main()
