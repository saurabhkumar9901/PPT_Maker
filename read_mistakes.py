from pptx import Presentation
import sys

def read_mistakes(filepath):
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"Error reading file: {e}")
        return
        
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join([p.text for p in shape.text_frame.paragraphs])
                if text.strip():
                    print(f"[{shape.shape_type}]: {text}")

if __name__ == "__main__":
    file_path = r"d:\PPT Maker\Common Mistakes and overall guide to improve slides.pptx"
    read_mistakes(file_path)
