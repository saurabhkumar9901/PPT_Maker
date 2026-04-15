from pptx import Presentation
import sys

prs = Presentation(sys.argv[1])
slide = prs.slides[11] # Slide 12 (0-indexed)
print("Slide 12 shapes:")
for s in slide.shapes:
    text = ""
    if s.has_text_frame:
        text = " ".join([p.text for p in s.text_frame.paragraphs])
    print(f"Type: {s.shape_type}, Name: {s.name}, W: {s.width}, H: {s.height}, Text: {text[:60]}")
