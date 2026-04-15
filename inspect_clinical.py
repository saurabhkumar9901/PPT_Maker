from pptx import Presentation
import sys

prs = Presentation(sys.argv[1])
for idx in [2, 9]: # Slide 3 and Slide 10 (0-indexed)
    print(f"\n==== SLIDE {idx+1} ====")
    slide = prs.slides[idx]
    for s in slide.shapes:
        text = ""
        font_size = "N/A"
        if s.has_text_frame:
            text = " ".join([p.text for p in s.text_frame.paragraphs])
            if s.text_frame.paragraphs and s.text_frame.paragraphs[0].runs:
                font_size = str(s.text_frame.paragraphs[0].runs[0].font.size)
        print(f"Type: {s.shape_type}, Name: {s.name}, W: {s.width}, H: {s.height}, Text: {text[:80]}, Size: {font_size}")
